from collections.abc import Sequence
from pptx import Presentation
from pptx.util import Inches


def read_questions_answers(file_path):
    questions_dict = {}

    with open(file_path, 'r') as file:
        current_question = None

        for line in file:
            line = line.strip()

            if line.startswith('Q'):
                # New question found
                current_question = line
                questions_dict[current_question] = []
            elif line.startswith('A'):
                # Answer for the current question
                if current_question:
                    questions_dict[current_question].append(line)

    return questions_dict







def add_slides(presentation, q_dict):

    # Add a slide for each question and answer
    for question, answers in q_dict.items():
        print(question)
        # Add another slide with a different layout
        bullet_slide_layout = presentation.slide_layouts[1]  # Use a bullet point layout
        slide2 = presentation.slides.add_slide(bullet_slide_layout)

        shapes = slide2.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = question

        for answer in answers:
            print(answer)
            # Add the answer as a bullet point
            content_frame = body_shape.text_frame
            content_frame.text = answer
            print()


# Create a presentation object
presentation = Presentation()

# Add a slide to the presentation
slide_layout = presentation.slide_layouts[0]  # Use the first slide layout (Title Slide)
slide = presentation.slides.add_slide(slide_layout)

# Add a title and content to the slide
title = slide.shapes.title
title.text = "IB Computer Science Revision"

content = slide.placeholders[1]  # Assuming content is in the second placeholder
content.text = "Questions and Answers."

# Replace 'your_file.txt' with the path to your actual text file
questions_answers_dict = read_questions_answers("Unit_1.txt")
add_slides(presentation, questions_answers_dict)
questions_answers_dict = read_questions_answers("Unit_2.txt")
add_slides(presentation, questions_answers_dict)
questions_answers_dict = read_questions_answers("Unit_3.txt")
add_slides(presentation, questions_answers_dict)
questions_answers_dict = read_questions_answers("Unit_4.txt")
add_slides(presentation, questions_answers_dict)


# Save the updated presentation
presentation.save("my_updated_presentation.pptx")
